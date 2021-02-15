import numpy as np
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches 
from io import StringIO
import os
import xlrd

wb = xlrd.open_workbook('test.xls')
sh1 = wb.sheet_by_name(u'Sheet1')

print(sh1.col_values(0))  # column 0
print(sh1.col_values(1)) # column 1
x = sh1.col_values(0)  # column 0
y = sh1.col_values(1)  # column 1
plt.plot(x, y)
plt.title('Unemployment Rate Vs Year')
plt.xlabel('Year')
plt.ylabel('Unemployment Rate')
fig1 = plt.gcf()
plt.show()
plt.draw()
fig1.savefig("image.png",dip=100)

prs = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[8])
placeholder = slide.placeholders[1]
picture = placeholder.insert_picture('image.png')
prs.save("graph.pptx")
os.startfile("graph.pptx")

